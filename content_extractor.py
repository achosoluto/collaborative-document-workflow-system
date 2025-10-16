"""
Advanced Document Content Extraction and Preprocessing Pipeline
Handles multiple document formats and prepares content for summarization and analysis
"""

import os
import re
import logging
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from datetime import datetime
import json

# Document processing libraries
import PyPDF2
from docx import Document
import pandas as pd
from bs4 import BeautifulSoup
import openpyxl
from pptx import Presentation

# Text processing and AI
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
import spacy
from transformers import pipeline
import torch

logger = logging.getLogger(__name__)


class DocumentContentExtractor:
    """Advanced document content extraction with preprocessing capabilities"""

    def __init__(self):
        self.supported_extensions = {
            '.pdf', '.docx', '.doc', '.txt', '.xlsx', '.xls', '.pptx', '.ppt',
            '.html', '.htm', '.json', '.csv', '.md', '.rtf'
        }

        # Initialize NLP components
        self._initialize_nlp()

        # Content processing settings
        self.max_chunk_size = 1000
        self.chunk_overlap = 100
        self.min_content_length = 50

    def _initialize_nlp(self):
        """Initialize NLP components"""
        try:
            # Download required NLTK data
            nltk_data = ['punkt', 'stopwords', 'wordnet', 'averaged_perceptron_tagger']
            for data in nltk_data:
                try:
                    nltk.data.find(f'tokenizers/punkt' if data == 'punkt' else
                                 f'corpora/{data}' if data != 'punkt' else f'tokenizers/punkt')
                except LookupError:
                    nltk.download(data, quiet=True)

            # Load spaCy model
            try:
                self.nlp = spacy.load('en_core_web_sm')
            except OSError:
                logger.info("Downloading spaCy model...")
                os.system('python -m spacy download en_core_web_sm')
                self.nlp = spacy.load('en_core_web_sm')

        except Exception as e:
            logger.error(f"Error initializing NLP components: {e}")
            self.nlp = None

    def extract_content(self, file_path: str) -> Dict[str, Any]:
        """
        Extract content from document with metadata

        Returns:
            Dictionary containing content, metadata, and processing information
        """
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                return self._error_response(f"File not found: {file_path}")

            # Get file extension and validate
            extension = file_path.suffix.lower()
            if extension not in self.supported_extensions:
                return self._error_response(f"Unsupported file type: {extension}")

            # Extract based on file type
            if extension == '.pdf':
                return self._extract_pdf_content(file_path)
            elif extension in ['.docx', '.doc']:
                return self._extract_docx_content(file_path)
            elif extension in ['.xlsx', '.xls']:
                return self._extract_excel_content(file_path)
            elif extension in ['.pptx', '.ppt']:
                return self._extract_powerpoint_content(file_path)
            elif extension in ['.html', '.htm']:
                return self._extract_html_content(file_path)
            elif extension == '.txt':
                return self._extract_text_content(file_path)
            elif extension == '.json':
                return self._extract_json_content(file_path)
            elif extension == '.csv':
                return self._extract_csv_content(file_path)
            elif extension == '.md':
                return self._extract_markdown_content(file_path)
            else:
                return self._error_response(f"Handler not implemented for: {extension}")

        except Exception as e:
            logger.error(f"Error extracting content from {file_path}: {e}")
            return self._error_response(f"Extraction error: {str(e)}")

    def _extract_pdf_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from PDF files"""
        try:
            content = []
            metadata = {}

            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)

                # Extract metadata
                if pdf_reader.metadata:
                    metadata = {
                        'title': pdf_reader.metadata.get('/Title', ''),
                        'author': pdf_reader.metadata.get('/Author', ''),
                        'subject': pdf_reader.metadata.get('/Subject', ''),
                        'creator': pdf_reader.metadata.get('/Creator', ''),
                        'producer': pdf_reader.metadata.get('/Producer', ''),
                        'creation_date': pdf_reader.metadata.get('/CreationDate', ''),
                        'pages': len(pdf_reader.pages)
                    }

                # Extract text from each page
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        text = page.extract_text()
                        if text.strip():
                            content.append({
                                'page': page_num + 1,
                                'content': text.strip(),
                                'char_count': len(text)
                            })
                    except Exception as e:
                        logger.warning(f"Error extracting page {page_num + 1}: {e}")

            # Combine all content
            full_text = '\n\n'.join([item['content'] for item in content])

            return {
                'success': True,
                'file_path': str(file_path),
                'file_type': 'pdf',
                'content': full_text,
                'content_by_page': content,
                'metadata': metadata,
                'total_pages': len(content),
                'total_characters': len(full_text),
                'extraction_method': 'pypdf2'
            }

        except Exception as e:
            return self._error_response(f"PDF extraction error: {str(e)}")

    def _extract_docx_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from Word documents"""
        try:
            doc = Document(file_path)
            content = []
            metadata = {}

            # Extract document properties
            if doc.core_properties:
                metadata = {
                    'title': doc.core_properties.title or '',
                    'author': doc.core_properties.author or '',
                    'subject': doc.core_properties.subject or '',
                    'keywords': doc.core_properties.keywords or '',
                    'created': str(doc.core_properties.created) if doc.core_properties.created else '',
                    'modified': str(doc.core_properties.modified) if doc.core_properties.modified else ''
                }

            # Extract content by paragraphs
            for para_num, paragraph in enumerate(doc.paragraphs):
                text = paragraph.text.strip()
                if text:
                    content.append({
                        'paragraph': para_num + 1,
                        'content': text,
                        'style': paragraph.style.name if paragraph.style else 'Normal',
                        'char_count': len(text)
                    })

            # Extract tables
            tables = []
            for table_num, table in enumerate(doc.tables):
                table_content = []
                for row_num, row in enumerate(table.rows):
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_content.append({
                        'row': row_num + 1,
                        'content': row_text
                    })
                tables.append({
                    'table': table_num + 1,
                    'rows': len(table.rows),
                    'columns': len(table.columns) if table.rows else 0,
                    'content': table_content
                })

            full_text = '\n\n'.join([item['content'] for item in content])

            return {
                'success': True,
                'file_path': str(file_path),
                'file_type': 'docx',
                'content': full_text,
                'content_by_paragraph': content,
                'tables': tables,
                'metadata': metadata,
                'total_paragraphs': len(content),
                'total_tables': len(tables),
                'total_characters': len(full_text),
                'extraction_method': 'python-docx'
            }

        except Exception as e:
            return self._error_response(f"Word document extraction error: {str(e)}")

    def _extract_excel_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from Excel files"""
        try:
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            content = []
            metadata = {}

            # Extract workbook properties
            if workbook.properties:
                metadata = {
                    'title': workbook.properties.title or '',
                    'author': workbook.properties.author or '',
                    'created': str(workbook.properties.created) if workbook.properties.created else '',
                    'modified': str(workbook.properties.modified) if workbook.properties.modified else ''
                }

            # Extract content from each sheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]

                sheet_content = {
                    'sheet_name': sheet_name,
                    'rows': [],
                    'headers': []
                }

                # Extract headers (first row)
                first_row = next(sheet.iter_rows(min_row=1, max_row=1, values_only=True))
                if first_row:
                    sheet_content['headers'] = [str(cell) if cell is not None else '' for cell in first_row]

                # Extract data rows
                for row_num, row in enumerate(sheet.iter_rows(min_row=2, values_only=True), 2):
                    row_data = [str(cell) if cell is not None else '' for cell in row]
                    if any(row_data):  # Only include non-empty rows
                        sheet_content['rows'].append({
                            'row': row_num,
                            'data': row_data
                        })

                content.append(sheet_content)

            # Create text representation
            text_parts = []
            for sheet in content:
                text_parts.append(f"Sheet: {sheet['sheet_name']}")
                if sheet['headers']:
                    text_parts.append(f"Headers: {', '.join(sheet['headers'])}")
                for row in sheet['rows'][:100]:  # Limit to first 100 rows for text
                    text_parts.append(f"Row {row['row']}: {', '.join(row['data'])}")
                text_parts.append("")  # Empty line between sheets

            full_text = '\n'.join(text_parts)

            return {
                'success': True,
                'file_path': str(file_path),
                'file_type': 'excel',
                'content': full_text,
                'sheets': content,
                'metadata': metadata,
                'total_sheets': len(content),
                'total_characters': len(full_text),
                'extraction_method': 'openpyxl'
            }

        except Exception as e:
            return self._error_response(f"Excel extraction error: {str(e)}")

    def _extract_text_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()

            # Basic metadata for text files
            lines = content.split('\n')
            metadata = {
                'lines': len(lines),
                'encoding': 'utf-8'
            }

            return {
                'success': True,
                'file_path': str(file_path),
                'file_type': 'text',
                'content': content,
                'metadata': metadata,
                'total_lines': len(lines),
                'total_characters': len(content),
                'extraction_method': 'text_reader'
            }

        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        content = file.read()
                    break
                except UnicodeDecodeError:
                    continue
            else:
                return self._error_response("Unable to decode text file with available encodings")

            return {
                'success': True,
                'file_path': str(file_path),
                'file_type': 'text',
                'content': content,
                'metadata': {'encoding': encoding, 'lines': len(content.split('\n'))},
                'total_characters': len(content),
                'extraction_method': f'text_reader_{encoding}'
            }

        except Exception as e:
            return self._error_response(f"Text file extraction error: {str(e)}")

    def _extract_html_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from HTML files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()

            soup = BeautifulSoup(html_content, 'html.parser')

            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()

            # Extract text content
            text = soup.get_text()

            # Extract metadata from HTML meta tags
            metadata = {}
            title_tag = soup.find('title')
            if title_tag:
                metadata['title'] = title_tag.get_text().strip()

            for meta in soup.find_all('meta'):
                name = meta.get('name') or meta.get('property')
                content = meta.get('content')
                if name and content:
                    metadata[name] = content

            # Clean up text
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            clean_text = '\n'.join(lines)

            return {
                'success': True,
                'file_path': str(file_path),
                'file_type': 'html',
                'content': clean_text,
                'raw_html': html_content,
                'metadata': metadata,
                'total_characters': len(clean_text),
                'extraction_method': 'beautifulsoup'
            }

        except Exception as e:
            return self._error_response(f"HTML extraction error: {str(e)}")

    def _extract_json_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from JSON files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)

            # Convert to text representation
            content = json.dumps(data, indent=2, ensure_ascii=False)

            return {
                'success': True,
                'file_path': str(file_path),
                'file_type': 'json',
                'content': content,
                'structured_data': data,
                'metadata': {'valid_json': True},
                'total_characters': len(content),
                'extraction_method': 'json_reader'
            }

        except json.JSONDecodeError as e:
            return self._error_response(f"Invalid JSON format: {str(e)}")
        except Exception as e:
            return self._error_response(f"JSON extraction error: {str(e)}")

    def _extract_csv_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from CSV files"""
        try:
            df = pd.read_csv(file_path)

            # Convert to text representation
            content = df.to_string(index=False)

            # Extract metadata
            metadata = {
                'columns': list(df.columns),
                'rows': len(df),
                'dtypes': {col: str(dtype) for col, dtype in df.dtypes.items()}
            }

            return {
                'success': True,
                'file_path': str(file_path),
                'file_type': 'csv',
                'content': content,
                'structured_data': df.to_dict('records'),
                'metadata': metadata,
                'total_rows': len(df),
                'total_columns': len(df.columns),
                'total_characters': len(content),
                'extraction_method': 'pandas'
            }

        except Exception as e:
            return self._error_response(f"CSV extraction error: {str(e)}")

    def _extract_markdown_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from Markdown files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()

            # Basic metadata extraction from frontmatter (if present)
            metadata = {}
            lines = content.split('\n')
            if lines and lines[0].strip() == '---':
                # Parse YAML frontmatter
                frontmatter_end = None
                for i, line in enumerate(lines[1:], 1):
                    if line.strip() == '---':
                        frontmatter_end = i
                        break

                if frontmatter_end:
                    import yaml
                    try:
                        frontmatter_text = '\n'.join(lines[1:frontmatter_end])
                        metadata = yaml.safe_load(frontmatter_text) or {}
                        content = '\n'.join(lines[frontmatter_end+1:])
                    except yaml.YAMLError:
                        pass  # Invalid YAML, use full content

            return {
                'success': True,
                'file_path': str(file_path),
                'file_type': 'markdown',
                'content': content,
                'metadata': metadata,
                'total_characters': len(content),
                'extraction_method': 'markdown_reader'
            }

        except Exception as e:
            return self._error_response(f"Markdown extraction error: {str(e)}")

    def _extract_powerpoint_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from PowerPoint files"""
        try:
            presentation = Presentation(file_path)
            content = []

            # Extract slide content
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = []

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)

                if slide_text:
                    content.append({
                        'slide': slide_num + 1,
                        'content': '\n'.join(slide_text),
                        'shape_count': len([s for s in slide.shapes if hasattr(s, "text")])
                    })

            full_text = '\n\n'.join([item['content'] for item in content])

            return {
                'success': True,
                'file_path': str(file_path),
                'file_type': 'powerpoint',
                'content': full_text,
                'slides': content,
                'metadata': {'total_slides': len(content)},
                'total_slides': len(content),
                'total_characters': len(full_text),
                'extraction_method': 'python-pptx'
            }

        except Exception as e:
            return self._error_response(f"PowerPoint extraction error: {str(e)}")

    def _error_response(self, message: str) -> Dict[str, Any]:
        """Return standardized error response"""
        return {
            'success': False,
            'error': message,
            'file_path': '',
            'file_type': 'unknown',
            'content': '',
            'metadata': {},
            'total_characters': 0,
            'extraction_method': 'error'
        }


class ContentPreprocessor:
    """Preprocesses extracted content for summarization and analysis"""

    def __init__(self):
        self.stop_words = set(stopwords.words('english'))
        self.lemmatizer = nltk.WordNetLemmatizer()

    def preprocess_content(self, extracted_content: Dict[str, Any]) -> Dict[str, Any]:
        """
        Preprocess extracted content for better analysis

        Args:
            extracted_content: Output from DocumentContentExtractor

        Returns:
            Enhanced content with preprocessing information
        """
        if not extracted_content['success']:
            return extracted_content

        content = extracted_content['content']
        if not content or len(content.strip()) < self.min_content_length:
            return {
                **extracted_content,
                'preprocessing': {
                    'status': 'insufficient_content',
                    'reason': 'Content too short for meaningful preprocessing'
                }
            }

        # Perform preprocessing steps
        preprocessing_result = {
            'status': 'completed',
            'steps_performed': [],
            'statistics': {}
        }

        # Step 1: Text cleaning
        cleaned_content = self._clean_text(content)
        preprocessing_result['steps_performed'].append('text_cleaning')

        # Step 2: Sentence tokenization
        sentences = sent_tokenize(cleaned_content)
        preprocessing_result['steps_performed'].append('sentence_tokenization')
        preprocessing_result['statistics']['sentence_count'] = len(sentences)

        # Step 3: Word tokenization and analysis
        words = word_tokenize(cleaned_content.lower())
        filtered_words = [w for w in words if w not in self.stop_words and w.isalnum()]
        preprocessing_result['steps_performed'].append('word_tokenization')
        preprocessing_result['statistics']['word_count'] = len(words)
        preprocessing_result['statistics']['filtered_word_count'] = len(filtered_words)

        # Step 4: Named entity recognition (if spaCy available)
        entities = []
        if self.nlp:
            try:
                doc = self.nlp(cleaned_content[:10000])  # Limit for performance
                entities = [(ent.text, ent.label_) for ent in doc.ents]
                preprocessing_result['steps_performed'].append('named_entity_recognition')
            except Exception as e:
                logger.warning(f"NER failed: {e}")

        preprocessing_result['statistics']['named_entities'] = len(entities)

        # Step 5: Keyword extraction
        keywords = self._extract_keywords(filtered_words, sentences)
        preprocessing_result['steps_performed'].append('keyword_extraction')
        preprocessing_result['statistics']['keywords_extracted'] = len(keywords)

        # Step 6: Content chunking for processing
        chunks = self._create_chunks(sentences)
        preprocessing_result['steps_performed'].append('content_chunking')
        preprocessing_result['statistics']['chunks_created'] = len(chunks)

        # Update extracted content with preprocessing results
        enhanced_content = {
            **extracted_content,
            'preprocessed_content': cleaned_content,
            'sentences': sentences,
            'filtered_words': filtered_words,
            'keywords': keywords,
            'named_entities': entities,
            'chunks': chunks,
            'preprocessing': preprocessing_result
        }

        return enhanced_content

    def _clean_text(self, text: str) -> str:
        """Clean and normalize text"""
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text)

        # Remove page numbers and headers/footers (common in PDFs)
        text = re.sub(r'\n\s*\d+\s*\n', '\n', text)  # Page numbers
        text = re.sub(r'^\s*\d+\s*$', '', text, flags=re.MULTILINE)  # Standalone numbers

        # Fix common OCR errors (if any)
        text = re.sub(r'(\w+)\s*-\s*\n\s*(\w+)', r'\1\2', text)  # Hyphenated words across lines

        return text.strip()

    def _extract_keywords(self, words: List[str], sentences: List[str]) -> List[Tuple[str, float]]:
        """Extract keywords with relevance scores"""
        from collections import Counter

        # Word frequency analysis
        word_freq = Counter(words)
        total_words = len(words)

        # Filter and score keywords
        keywords = []
        for word, freq in word_freq.most_common(50):  # Top 50 candidates
            if len(word) > 3 and freq > 1:  # Meaningful words that appear more than once
                # Calculate relevance score
                score = (freq / total_words) * (1 + len(word) / 10)  # Longer words get slight boost
                keywords.append((word, score))

        # Sort by score and return top 20
        keywords.sort(key=lambda x: x[1], reverse=True)
        return keywords[:20]

    def _create_chunks(self, sentences: List[str]) -> List[Dict[str, Any]]:
        """Create overlapping chunks for processing"""
        chunks = []
        current_chunk = []
        current_length = 0

        for i, sentence in enumerate(sentences):
            sentence_length = len(sentence)

            # Check if adding this sentence would exceed max chunk size
            if current_length + sentence_length > self.max_chunk_size and current_chunk:
                # Create chunk from current sentences
                chunk_text = ' '.join(current_chunk)
                chunks.append({
                    'chunk_id': f"chunk_{len(chunks)}",
                    'content': chunk_text,
                    'sentence_count': len(current_chunk),
                    'start_sentence': i - len(current_chunk),
                    'end_sentence': i - 1,
                    'char_count': current_length
                })

                # Start new chunk with overlap
                overlap_sentences = current_chunk[-1:] if len(current_chunk) > 0 else []
                current_chunk = overlap_sentences + [sentence]
                current_length = sum(len(s) for s in current_chunk)
            else:
                current_chunk.append(sentence)
                current_length += sentence_length

        # Add final chunk
        if current_chunk:
            chunk_text = ' '.join(current_chunk)
            chunks.append({
                'chunk_id': f"chunk_{len(chunks)}",
                'content': chunk_text,
                'sentence_count': len(current_chunk),
                'start_sentence': len(sentences) - len(current_chunk),
                'end_sentence': len(sentences) - 1,
                'char_count': current_length
            })

        return chunks

    @property
    def min_content_length(self) -> int:
        """Minimum content length for preprocessing"""
        return getattr(self, '_min_content_length', 50)

    @min_content_length.setter
    def min_content_length(self, value: int):
        """Set minimum content length"""
        self._min_content_length = value


class ContentProcessingPipeline:
    """Complete content processing pipeline"""

    def __init__(self):
        self.extractor = DocumentContentExtractor()
        self.preprocessor = ContentPreprocessor()

    def process_document(self, file_path: str) -> Dict[str, Any]:
        """
        Complete document processing pipeline

        Args:
            file_path: Path to document file

        Returns:
            Fully processed document with extracted and preprocessed content
        """
        # Step 1: Extract content
        extracted = self.extractor.extract_content(file_path)

        if not extracted['success']:
            return extracted

        # Step 2: Preprocess content
        processed = self.preprocessor.preprocess_content(extracted)

        # Add pipeline metadata
        processed['pipeline_info'] = {
            'extraction_timestamp': datetime.now().isoformat(),
            'pipeline_version': '1.0',
            'processing_steps': ['extraction', 'preprocessing'],
            'file_path': file_path
        }

        return processed

    def process_multiple_documents(self, file_paths: List[str]) -> List[Dict[str, Any]]:
        """Process multiple documents in batch"""
        results = []

        for file_path in file_paths:
            try:
                result = self.process_document(file_path)
                results.append(result)
            except Exception as e:
                logger.error(f"Error processing {file_path}: {e}")
                results.append({
                    'success': False,
                    'error': str(e),
                    'file_path': file_path
                })

        return results